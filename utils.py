import os
import re
import fitz  # PDF
import docx  # Word
import pptx  # PowerPoint
import streamlit as st

# === Universal Text Extractor ===
def extract_text_from_file(path):
    ext = os.path.splitext(path)[1].lower()

    try:
        if ext == ".pdf":
            doc = fitz.open(path)
            return " ".join([page.get_text() for page in doc])

        elif ext == ".txt":
            with open(path, "r", encoding="utf-8") as f:
                return f.read()

        elif ext == ".docx":
            document = docx.Document(path)
            return "\n".join([para.text for para in document.paragraphs])

        elif ext == ".pptx":
            presentation = pptx.Presentation(path)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text

        elif ext == ".md":
            with open(path, "r", encoding="utf-8") as f:
                return f.read()

        else:
            return ""

    except Exception as e:
        st.error(f"❌ Failed to extract text from {ext.upper()} file: {e}")
        return ""

# === Chunking ===
def chunk_text(text, chunk_size=200):
    sentences = re.split(r'(?<=[.!?]) +', text)
    chunks, chunk, count = [], [], 0
    for s in sentences:
        words = s.split()
        if count + len(words) > chunk_size:
            chunks.append(" ".join(chunk))
            chunk, count = [], 0
        chunk += words
        count += len(words)
    if chunk:
        chunks.append(" ".join(chunk))
    return chunks

# === Stats ===
def get_stats():
    try:
        from config import collection
        # Only get chunks related to the currently uploaded PDFs
        current_pdfs = st.session_state.get("pdf_names", [])
        all_data = collection.get()
        chunk_count = sum(1 for meta in all_data["metadatas"] if meta["pdf"] in current_pdfs)
        question_count = len(st.session_state.get("chat_history", []))
        return {"chunks": chunk_count, "questions": question_count}
    except Exception as e:
        st.warning(f"⚠️ Failed to get stats from ChromaDB: {e}")
        return {"chunks": 0, "questions": 0}
